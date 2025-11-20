"""Loaders for Microsoft Office documents (Word, Excel, PowerPoint)"""
import os
from typing import List, Optional, Union, BinaryIO
import io
import logging
from .base import BaseDocumentLoader, LoadedDocument, DocumentType, DocumentChunk

logger = logging.getLogger(__name__)


class WordLoader(BaseDocumentLoader):
    """Loader for Microsoft Word documents (.doc, .docx)"""

    def supported_extensions(self) -> List[str]:
        """Return supported file extensions"""
        return [".doc", ".docx"]

    def load(
        self,
        file_path: Optional[str] = None,
        file_content: Optional[Union[bytes, BinaryIO]] = None,
        **kwargs
    ) -> LoadedDocument:
        """Load Word document"""
        if not file_path and not file_content:
            raise ValueError("Either file_path or file_content must be provided")

        # Get file content
        if file_path:
            file_size = os.path.getsize(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
        else:
            if isinstance(file_content, bytes):
                doc_bytes = file_content
            else:
                doc_bytes = file_content.read()
            file_size = len(doc_bytes)
            file_ext = ".docx"  # Default assumption

        self.validate_file_size(file_size)

        # Extract text based on format
        if file_ext == ".docx":
            text_content, page_count = self._extract_docx(file_path, file_content)
        else:  # .doc
            text_content, page_count = self._extract_doc(file_path, file_content)

        metadata = self.extract_metadata(file_path)
        metadata.update({
            "document_format": file_ext,
            "page_count": page_count
        })

        # Create chunks
        chunks = None
        if self.enable_chunking and text_content:
            from ..chunking.strategies import get_chunking_strategy
            chunker = get_chunking_strategy(
                strategy_name="semantic",
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )
            text_chunks = chunker.chunk(text_content, metadata)

            chunks = [
                DocumentChunk(
                    content=chunk.content,
                    chunk_index=chunk.index,
                    total_chunks=chunk.total_chunks,
                    metadata=chunk.metadata
                )
                for chunk in text_chunks
            ]

        return LoadedDocument(
            content=text_content,
            document_type=DocumentType.WORD,
            file_size=file_size,
            page_count=page_count,
            metadata=metadata,
            chunks=chunks
        )

    def _extract_docx(
        self,
        file_path: Optional[str],
        file_content: Optional[Union[bytes, BinaryIO]]
    ) -> tuple[str, int]:
        """Extract text from .docx file"""
        try:
            import docx

            if file_path:
                doc = docx.Document(file_path)
            else:
                if isinstance(file_content, bytes):
                    doc = docx.Document(io.BytesIO(file_content))
                else:
                    doc = docx.Document(file_content)

            # Extract paragraphs
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Extract tables
            for table in doc.tables:
                table_text = self._extract_table_text(table)
                if table_text:
                    text_parts.append(table_text)

            full_text = "\n\n".join(text_parts)

            # Estimate page count (rough approximation)
            page_count = max(1, len(full_text) // 3000)

            return full_text, page_count

        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {e}")
            return "", 0

    def _extract_doc(
        self,
        file_path: Optional[str],
        file_content: Optional[Union[bytes, BinaryIO]]
    ) -> tuple[str, int]:
        """Extract text from .doc file (older format)"""
        try:
            import textract

            if file_path:
                text = textract.process(file_path).decode('utf-8')
            else:
                # textract requires file path, save temporarily
                import tempfile
                with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as tmp:
                    if isinstance(file_content, bytes):
                        tmp.write(file_content)
                    else:
                        tmp.write(file_content.read())
                    tmp_path = tmp.name

                text = textract.process(tmp_path).decode('utf-8')
                os.unlink(tmp_path)

            page_count = max(1, len(text) // 3000)
            return text, page_count

        except ImportError:
            logger.error("textract not available for .doc files - install with: pip install textract")
            return "", 0
        except Exception as e:
            logger.error(f"Failed to extract text from DOC: {e}")
            return "", 0

    def _extract_table_text(self, table) -> str:
        """Extract text from Word table"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)


class ExcelLoader(BaseDocumentLoader):
    """Loader for Microsoft Excel documents (.xls, .xlsx)"""

    def supported_extensions(self) -> List[str]:
        """Return supported file extensions"""
        return [".xls", ".xlsx", ".xlsm"]

    def load(
        self,
        file_path: Optional[str] = None,
        file_content: Optional[Union[bytes, BinaryIO]] = None,
        **kwargs
    ) -> LoadedDocument:
        """Load Excel document"""
        if not file_path and not file_content:
            raise ValueError("Either file_path or file_content must be provided")

        # Get file size
        if file_path:
            file_size = os.path.getsize(file_path)
        else:
            if isinstance(file_content, bytes):
                file_size = len(file_content)
            else:
                file_size = len(file_content.read())
                file_content.seek(0)

        self.validate_file_size(file_size)

        # Extract text
        text_content, sheet_count = self._extract_excel(file_path, file_content)

        metadata = self.extract_metadata(file_path)
        metadata.update({
            "sheet_count": sheet_count,
            "page_count": sheet_count  # Treat sheets as pages
        })

        # Create chunks
        chunks = None
        if self.enable_chunking and text_content:
            from ..chunking.strategies import get_chunking_strategy
            chunker = get_chunking_strategy(
                strategy_name="fixed_size",
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )
            text_chunks = chunker.chunk(text_content, metadata)

            chunks = [
                DocumentChunk(
                    content=chunk.content,
                    chunk_index=chunk.index,
                    total_chunks=chunk.total_chunks,
                    metadata=chunk.metadata
                )
                for chunk in text_chunks
            ]

        return LoadedDocument(
            content=text_content,
            document_type=DocumentType.EXCEL,
            file_size=file_size,
            page_count=sheet_count,
            metadata=metadata,
            chunks=chunks
        )

    def _extract_excel(
        self,
        file_path: Optional[str],
        file_content: Optional[Union[bytes, BinaryIO]]
    ) -> tuple[str, int]:
        """Extract text from Excel file"""
        try:
            import openpyxl

            if file_path:
                workbook = openpyxl.load_workbook(file_path, data_only=True)
            else:
                if isinstance(file_content, bytes):
                    workbook = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
                else:
                    workbook = openpyxl.load_workbook(file_content, data_only=True)

            text_parts = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"=== Sheet: {sheet_name} ===\n"]

                for row in sheet.iter_rows(values_only=True):
                    # Filter out None values and convert to strings
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                if len(sheet_text) > 1:  # Has content beyond header
                    text_parts.append("\n".join(sheet_text))

            full_text = "\n\n".join(text_parts)
            sheet_count = len(workbook.sheetnames)

            return full_text, sheet_count

        except Exception as e:
            logger.error(f"Failed to extract text from Excel: {e}")
            return "", 0


class PowerPointLoader(BaseDocumentLoader):
    """Loader for Microsoft PowerPoint documents (.ppt, .pptx)"""

    def supported_extensions(self) -> List[str]:
        """Return supported file extensions"""
        return [".ppt", ".pptx"]

    def load(
        self,
        file_path: Optional[str] = None,
        file_content: Optional[Union[bytes, BinaryIO]] = None,
        **kwargs
    ) -> LoadedDocument:
        """Load PowerPoint document"""
        if not file_path and not file_content:
            raise ValueError("Either file_path or file_content must be provided")

        # Get file size
        if file_path:
            file_size = os.path.getsize(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
        else:
            if isinstance(file_content, bytes):
                file_size = len(file_content)
            else:
                ppt_bytes = file_content.read()
                file_content = io.BytesIO(ppt_bytes)
                file_size = len(ppt_bytes)
            file_ext = ".pptx"

        self.validate_file_size(file_size)

        # Extract text
        if file_ext == ".pptx":
            text_content, slide_count = self._extract_pptx(file_path, file_content)
        else:
            text_content, slide_count = self._extract_ppt(file_path, file_content)

        metadata = self.extract_metadata(file_path)
        metadata.update({
            "slide_count": slide_count,
            "page_count": slide_count
        })

        # Create chunks
        chunks = None
        if self.enable_chunking and text_content:
            from ..chunking.strategies import get_chunking_strategy
            chunker = get_chunking_strategy(
                strategy_name="semantic",
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap
            )
            text_chunks = chunker.chunk(text_content, metadata)

            chunks = [
                DocumentChunk(
                    content=chunk.content,
                    chunk_index=chunk.index,
                    total_chunks=chunk.total_chunks,
                    metadata=chunk.metadata
                )
                for chunk in text_chunks
            ]

        return LoadedDocument(
            content=text_content,
            document_type=DocumentType.POWERPOINT,
            file_size=file_size,
            page_count=slide_count,
            metadata=metadata,
            chunks=chunks
        )

    def _extract_pptx(
        self,
        file_path: Optional[str],
        file_content: Optional[Union[bytes, BinaryIO]]
    ) -> tuple[str, int]:
        """Extract text from .pptx file"""
        try:
            from pptx import Presentation

            if file_path:
                prs = Presentation(file_path)
            else:
                prs = Presentation(file_content)

            text_parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== Slide {i} ===\n"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))

            full_text = "\n\n".join(text_parts)
            slide_count = len(prs.slides)

            return full_text, slide_count

        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {e}")
            return "", 0

    def _extract_ppt(
        self,
        file_path: Optional[str],
        file_content: Optional[Union[bytes, BinaryIO]]
    ) -> tuple[str, int]:
        """Extract text from .ppt file (older format)"""
        try:
            import textract

            if file_path:
                text = textract.process(file_path).decode('utf-8')
            else:
                # Save temporarily
                import tempfile
                with tempfile.NamedTemporaryFile(delete=False, suffix='.ppt') as tmp:
                    if isinstance(file_content, bytes):
                        tmp.write(file_content)
                    else:
                        tmp.write(file_content.read())
                    tmp_path = tmp.name

                text = textract.process(tmp_path).decode('utf-8')
                os.unlink(tmp_path)

            # Estimate slide count
            slide_count = text.count("Slide ") or max(1, len(text) // 500)
            return text, slide_count

        except ImportError:
            logger.error("textract not available for .ppt files")
            return "", 0
        except Exception as e:
            logger.error(f"Failed to extract text from PPT: {e}")
            return "", 0
